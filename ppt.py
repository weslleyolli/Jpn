from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paletas ─────────────────────────────────────────────────────────────────
POKEMON = {
    "bg":     RGBColor(0xFF, 0xD7, 0x00),
    "accent": RGBColor(0xCC, 0x00, 0x00),
    "dark":   RGBColor(0x1A, 0x1A, 0x2E),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xFF, 0xF0, 0x80),
}

NARUTO = {
    "bg":     RGBColor(0xFF, 0x6B, 0x00),
    "accent": RGBColor(0xFF, 0xC3, 0x00),
    "dark":   RGBColor(0x12, 0x07, 0x00),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xFF, 0xE0, 0xB2),
}

JJK = {
    "bg":     RGBColor(0x0D, 0x00, 0x1A),
    "accent": RGBColor(0x8A, 0x2B, 0xE2),
    "dark":   RGBColor(0x05, 0x00, 0x0D),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xCE, 0x93, 0xD8),
}

TOKYO = {
    "bg":     RGBColor(0x0A, 0x0A, 0x2E),
    "accent": RGBColor(0xE0, 0x30, 0x70),
    "dark":   RGBColor(0x05, 0x05, 0x18),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xB0, 0xC4, 0xDE),
}

DISNEY = {
    "bg":     RGBColor(0x00, 0x3D, 0x99),
    "accent": RGBColor(0xFF, 0xD7, 0x00),
    "dark":   RGBColor(0x00, 0x1A, 0x4D),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xAD, 0xD8, 0xE6),
}

FUJI = {
    "bg":     RGBColor(0x00, 0x2B, 0x5C),
    "accent": RGBColor(0xFF, 0x45, 0x00),
    "dark":   RGBColor(0x00, 0x10, 0x24),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xB0, 0xD0, 0xFF),
}

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────────────

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def circle(slide, cx, cy, r, color):
    sh = slide.shapes.add_shape(9, cx - r, cy - r, r * 2, r * 2)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def text(slide, txt, x, y, w, h, size=20, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return tb

def footer(slide, msg, bg_color, txt_color):
    rect(slide, 0, H - Inches(0.3), W, Inches(0.3), bg_color)
    text(slide, msg, Inches(0.3), H - Inches(0.3), W - Inches(0.6),
         Inches(0.3), size=11, bold=True, color=txt_color, align=PP_ALIGN.CENTER)

# ── Modelos de slide por tema ─────────────────────────────────────────────────

def make_naruto_slide(data, titulo, subtitulo, quote):
    """data = list of (icon, atividade, detalhe)"""
    s = blank()
    bg(s, NARUTO["dark"])
    for i in range(14):
        c = RGBColor(0x12 + i * 3, 0x07 + i, 0x00)
        rect(s, 0, Inches(i * 7.5 / 14), W, Inches(7.5 / 14 + 0.05), c)

    rect(s, 0, 0, W, Inches(1.45), NARUTO["bg"])
    rect(s, 0, Inches(1.45), W, Inches(0.07), NARUTO["accent"])
    rect(s, 0, 0, Inches(0.35), H, NARUTO["dark"])

    # Espiral
    circle(s, W - Inches(1.15), Inches(0.72), Inches(0.52), NARUTO["accent"])
    circle(s, W - Inches(1.15), Inches(0.72), Inches(0.34), RGBColor(0xFF, 0xA0, 0x30))
    circle(s, W - Inches(1.15), Inches(0.72), Inches(0.13), NARUTO["dark"])

    text(s, titulo, Inches(0.55), Inches(0.06), Inches(10.5), Inches(0.65),
         size=27, bold=True, color=NARUTO["dark"])
    text(s, subtitulo, Inches(0.55), Inches(0.72), Inches(10.5), Inches(0.65),
         size=19, bold=True, color=NARUTO["white"])

    for i, (icon, ativ, det) in enumerate(data):
        cy = Inches(1.68 + i * 0.92)
        rect(s, Inches(0.42), cy + Inches(0.08), Inches(0.1), Inches(0.62), NARUTO["bg"])
        text(s, icon + "  " + ativ, Inches(0.65), cy, Inches(10.8), Inches(0.5),
             size=16, bold=True, color=NARUTO["white"])
        text(s, det, Inches(0.65), cy + Inches(0.42), Inches(11.2), Inches(0.42),
             size=13, color=NARUTO["light"])

    footer(s, quote, NARUTO["bg"], NARUTO["dark"])
    return s


def make_pokemon_slide(data, titulo, subtitulo, quote):
    s = blank()
    bg(s, RGBColor(0x1A, 0x1A, 0x2E))

    rect(s, 0, 0, W, Inches(1.45), POKEMON["accent"])
    rect(s, 0, Inches(1.45), W, Inches(0.07), POKEMON["bg"])

    # Pokébola
    px, py = W - Inches(1.1), Inches(0.72)
    circle(s, px, py, Inches(0.58), POKEMON["white"])
    circle(s, px, py, Inches(0.56), POKEMON["accent"])
    rect(s, px - Inches(0.56), py - Inches(0.04), Inches(1.12), Inches(0.08), POKEMON["white"])
    circle(s, px, py, Inches(0.18), POKEMON["white"])
    circle(s, px, py, Inches(0.12), RGBColor(0xCC, 0xCC, 0xCC))

    text(s, titulo, Inches(0.4), Inches(0.07), Inches(10.5), Inches(0.65),
         size=27, bold=True, color=POKEMON["dark"])
    text(s, subtitulo, Inches(0.4), Inches(0.73), Inches(10.5), Inches(0.65),
         size=19, bold=True, color=POKEMON["white"])

    for i, (icon, ativ, det) in enumerate(data):
        cy = Inches(1.68 + i * 0.92)
        circle(s, Inches(0.57), cy + Inches(0.26), Inches(0.12), POKEMON["bg"])
        text(s, icon + "  " + ativ, Inches(0.78), cy, Inches(11.0), Inches(0.5),
             size=16, bold=True, color=POKEMON["bg"])
        text(s, det, Inches(0.78), cy + Inches(0.42), Inches(11.2), Inches(0.42),
             size=13, color=POKEMON["light"])

    footer(s, quote, POKEMON["accent"], POKEMON["dark"])
    return s


def make_jjk_slide(data, titulo, subtitulo, quote):
    s = blank()
    bg(s, JJK["bg"])
    for i in range(6):
        c = RGBColor(0x0D + i * 6, 0x00, 0x1A + i * 10)
        rect(s, Inches(i * 2.25), 0, Inches(2.22), H, c)

    rect(s, 0, 0, W, Inches(1.45), JJK["dark"])
    rect(s, 0, Inches(1.45), W, Inches(0.07), JJK["accent"])
    rect(s, 0, 0, Inches(0.4), H, JJK["accent"])

    circle(s, W - Inches(1.1), Inches(0.72), Inches(0.52), JJK["accent"])
    circle(s, W - Inches(1.1), Inches(0.72), Inches(0.36), JJK["dark"])
    circle(s, W - Inches(1.1), Inches(0.72), Inches(0.17), JJK["accent"])

    text(s, titulo, Inches(0.62), Inches(0.07), Inches(10.5), Inches(0.65),
         size=27, bold=True, color=JJK["accent"])
    text(s, subtitulo, Inches(0.62), Inches(0.73), Inches(10.5), Inches(0.65),
         size=19, bold=True, color=JJK["white"])

    for i, (icon, ativ, det) in enumerate(data):
        cy = Inches(1.68 + i * 0.92)
        rect(s, Inches(0.52), cy + Inches(0.1), Inches(0.07), Inches(0.52), JJK["accent"])
        text(s, icon + "  " + ativ, Inches(0.72), cy, Inches(11.0), Inches(0.5),
             size=16, bold=True, color=JJK["white"])
        text(s, det, Inches(0.72), cy + Inches(0.42), Inches(11.2), Inches(0.42),
             size=13, color=JJK["light"])

    footer(s, quote, JJK["accent"], JJK["dark"])
    return s


def make_disney_slide(data, titulo, subtitulo, quote, p=DISNEY):
    s = blank()
    bg(s, p["dark"])
    br, bg_, bb = p["bg"][0], p["bg"][1], p["bg"][2]
    for i in range(10):
        c = RGBColor(
            max(0, br   - i * 4),
            max(0, bg_  - i * 3),
            min(255, bb + i * 2),
        )
        rect(s, 0, Inches(i * 7.5 / 10), W, Inches(7.5 / 10 + 0.05), c)

    rect(s, 0, 0, W, Inches(1.45), p["bg"])
    rect(s, 0, Inches(1.45), W, Inches(0.07), p["accent"])

    # Estrela Disney
    for ang_off in range(4):
        cx_ = W - Inches(1.1 + ang_off * 0.08)
        cy_ = Inches(0.72)
        circle(s, cx_, cy_, Inches(0.08), p["accent"])
    circle(s, W - Inches(1.1), Inches(0.72), Inches(0.48), p["accent"])
    circle(s, W - Inches(1.1), Inches(0.72), Inches(0.3), p["dark"])
    circle(s, W - Inches(1.1), Inches(0.72), Inches(0.12), p["accent"])

    text(s, titulo, Inches(0.5), Inches(0.07), Inches(10.5), Inches(0.65),
         size=27, bold=True, color=p["accent"])
    text(s, subtitulo, Inches(0.5), Inches(0.73), Inches(10.5), Inches(0.65),
         size=19, bold=True, color=p["white"])

    for i, (icon, ativ, det) in enumerate(data):
        cy = Inches(1.68 + i * 0.92)
        circle(s, Inches(0.57), cy + Inches(0.26), Inches(0.14), p["accent"])
        text(s, icon + "  " + ativ, Inches(0.82), cy, Inches(11.0), Inches(0.5),
             size=16, bold=True, color=p["accent"])
        text(s, det, Inches(0.82), cy + Inches(0.42), Inches(11.2), Inches(0.42),
             size=13, color=p["light"])

    footer(s, quote, p["accent"], p["dark"])
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════════
s = blank()
bg(s, TOKYO["dark"])
for i in range(9):
    c = RGBColor(0x0A + i * 3, 0x0A + i * 2, 0x2E + i * 5)
    rect(s, 0, Inches(i * 7.5 / 9), W, Inches(7.5 / 9 + 0.05), c)

rect(s, 0, 0, Inches(0.42), H, TOKYO["accent"])
rect(s, Inches(0.42), 0, Inches(0.15), H, RGBColor(0xFF, 0xD7, 0x00))
rect(s, Inches(0.57), 0, Inches(0.15), H, NARUTO["bg"])
rect(s, Inches(0.72), 0, Inches(0.15), H, JJK["accent"])

# Círculos temáticos no canto
circle(s, W - Inches(1.4), Inches(1.3), Inches(0.9), POKEMON["accent"])
circle(s, W - Inches(1.4), Inches(1.3), Inches(0.65), RGBColor(0x1A, 0x1A, 0x2E))
circle(s, W - Inches(1.4), Inches(1.3), Inches(0.38), POKEMON["bg"])
circle(s, W - Inches(1.4), Inches(1.3), Inches(0.14), POKEMON["accent"])

circle(s, W - Inches(2.8), Inches(6.4), Inches(0.55), NARUTO["bg"])
circle(s, W - Inches(2.8), Inches(6.4), Inches(0.35), RGBColor(0xFF, 0xA0, 0x30))
circle(s, W - Inches(2.8), Inches(6.4), Inches(0.13), NARUTO["dark"])

circle(s, Inches(1.3), Inches(6.3), Inches(0.45), JJK["accent"])
circle(s, Inches(1.3), Inches(6.3), Inches(0.28), JJK["dark"])

text(s, "ROTEIRO DE VIAGEM", Inches(1.1), Inches(1.1), Inches(10), Inches(1.0),
     size=46, bold=True, color=TOKYO["light"])
text(s, "TOKYO  🗼", Inches(1.1), Inches(2.0), Inches(10), Inches(1.6),
     size=82, bold=True, color=TOKYO["accent"])
text(s, "22 de Abril  →  30 de Abril  •  9 dias de aventura",
     Inches(1.1), Inches(3.75), Inches(9.5), Inches(0.65),
     size=21, italic=True, color=TOKYO["light"])

# Badge dos animes
for i, (label, cor) in enumerate([("⚡ Pokemon", POKEMON["bg"]),
                                    ("🍥 Naruto",  NARUTO["bg"]),
                                    ("🩸 Jujutsu Kaisen", JJK["accent"])]):
    bx = Inches(1.1 + i * 3.55)
    rect(s, bx, Inches(4.65), Inches(3.35), Inches(0.58), cor)
    text(s, label, bx + Inches(0.1), Inches(4.67), Inches(3.15), Inches(0.52),
         size=16, bold=True, color=RGBColor(0x10, 0x10, 0x10), align=PP_ALIGN.CENTER)

rect(s, 0, H - Inches(0.35), W, Inches(0.35), TOKYO["accent"])
text(s, "Japan  •  Abril/2025  •  Uma aventura digna de protagonista de anime",
     Inches(0.3), H - Inches(0.35), W - Inches(0.6), Inches(0.35),
     size=12, bold=True, color=TOKYO["dark"], align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — VISÃO GERAL DO ROTEIRO
# ══════════════════════════════════════════════════════════════════════════════
s = blank()
bg(s, POKEMON["dark"])
rect(s, 0, 0, W, Inches(1.3), POKEMON["accent"])
rect(s, 0, Inches(1.3), W, Inches(0.07), POKEMON["bg"])

text(s, "⚡  MAPA DA AVENTURA  —  9 DIAS", 0, Inches(0.2), W, Inches(0.9),
     size=36, bold=True, color=POKEMON["dark"], align=PP_ALIGN.CENTER)

dias = [
    ("22/04", "🏯 Senso-ji • Ueno • Asakusa • Sky Tree • Ikebukuro",  NARUTO["bg"]),
    ("23/04", "⚡ Pokemon Center • Harajuku • Takeshita • Kiddy Land • Shinjuku", POKEMON["accent"]),
    ("24/04", "🤖 Robot Café • Tokyo Tower • Ginza • TeamLab Planets", JJK["accent"]),
    ("25/04", "🛍️ Shibuya • Capcom/Nintendo/Pokemon • Shibuya Sky",   POKEMON["bg"]),
    ("26/04", "🚄 BATE-VOLTA KYOTO",                                  RGBColor(0x2E, 0x7D, 0x32)),
    ("27/04", "🏰 TOKYO DISNEYLAND — do opening ao fechamento",        DISNEY["bg"]),
    ("28/04", "🌊 TOKYO DISNEYSEA — do opening ao fechamento",         RGBColor(0x00, 0x5F, 0x99)),
    ("29/04", "🎮 Akihabara • Rolê livre em grupo",                   NARUTO["accent"]),
    ("30/04", "🗻 Monte Fuji • Fuji-Q Highland  →  Osaka de Shinkansen", FUJI["accent"]),
]

for i, (dia, resumo, cor) in enumerate(dias):
    cy = Inches(1.52 + i * 0.655)
    rect(s, Inches(0.3), cy, Inches(1.35), Inches(0.58), cor)
    text(s, dia, Inches(0.3), cy + Inches(0.06), Inches(1.35), Inches(0.52),
         size=15, bold=True, color=RGBColor(0x10, 0x10, 0x10), align=PP_ALIGN.CENTER)
    rect(s, Inches(1.68), cy, Inches(10.8), Inches(0.58),
         RGBColor(0x25, 0x25, 0x3D))
    text(s, resumo, Inches(1.82), cy + Inches(0.06), Inches(10.5), Inches(0.5),
         size=14, color=POKEMON["light"])

footer(s, "Gotta go 'em all places!  •  9 dias de pura aventura japonesa",
       POKEMON["accent"], POKEMON["dark"])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 22/04: Senso-ji • Ueno • Asakusa • Sky Tree • Ikebukuro
# ══════════════════════════════════════════════════════════════════════════════
make_naruto_slide(
    [
        ("🏯", "Templo Senso-ji",   "Templo budista mais antigo de Tokyo — chegue cedo para evitar filas"),
        ("👀", "Asakusa",            "Bairro histórico ao redor do templo: ruelas, quimonos e artesanato"),
        ("🌳", "Ueno Park",          "Parque enorme com museus, zoológico e natureza no coração da cidade"),
        ("🗼", "Tokyo Sky Tree",     "634m de altura — a torre mais alta do Japão com vista 360°"),
        ("🎨", "Ikebukuro",          "Bairro animado: Sunshine City, lojas de anime e vida local"),
        ("🍣", "Kura Sushi ★",       "Sushi giratório com preços ótimos — experiência clássica japonesa"),
    ],
    "22/04 (TER)  —  SENSO-JI • UENO • ASAKUSA • SKY TREE • IKEBUKURO",
    "Primeira imersão na Tokyo tradicional e moderna!",
    "Dattebayo!  —  Primeiro dia de missão ninja em Tokyo!",
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 23/04: Pokemon Center • Harajuku • Shinjuku
# ══════════════════════════════════════════════════════════════════════════════
make_pokemon_slide(
    [
        ("⚡", "Pokemon Center Mega Tokyo",  "Loja oficial com exclusivos, pelúcias, TCG e colecionáveis raros"),
        ("🎀", "Harajuku",                   "Moda alternativa, cultura kawaii e cosplay nas ruas"),
        ("👕", "Takeshita Street",            "A rua mais colorida de Tokyo — crepes, lojas kawaii e moda única"),
        ("🧸", "Kiddy Land",                  "5 andares de brinquedos, figures, Hello Kitty e personagens"),
        ("✨", "Shinjuku",                    "Bairro que não dorme: luzes de néon, Golden Gai e Kabukicho"),
    ],
    "23/04 (QUA)  —  POKEMON CENTER • HARAJUKU • SHINJUKU",
    "Gotta catch 'em all — figurinhas, moda e néon!",
    "Você é o treinador mais preparado de Tokyo!",
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 24/04: Robot Café • Tokyo Tower • Ginza • TeamLab
# ══════════════════════════════════════════════════════════════════════════════
make_jjk_slide(
    [
        ("🤖", "Dawn Robot Café",        "Café operado por robôs — tecnologia japonesa do futuro num copo"),
        ("🗼", "Tokyo Tower",            "Ícone vermelho e branco de Tokyo — vista clássica da cidade"),
        ("🛍️", "Ginza",                  "O bairro mais sofisticado: lojas de grife, galerias e arquitetura"),
        ("🎨", "TeamLab Planets",        "Arte digital imersiva — entra literalmente dentro das obras"),
    ],
    "24/04 (QUI)  —  ROBOT CAFÉ • TOKYO TOWER • GINZA • TEAMLAB",
    "Energia amaldicoada de arte, tecnologia e neon!",
    '"Domínio Expandido" — TeamLab ativa sua melhor técnica proibida!',
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 25/04: Shibuya • Shopping • Shibuya Sky
# ══════════════════════════════════════════════════════════════════════════════
make_pokemon_slide(
    [
        ("🛍️", "Shibuya — Mecca das compras",    "O bairro mais movimentado do mundo — cada canto é fotogênico"),
        ("🎮", "Capcom Store",                    "Loja oficial da Capcom: Street Fighter, Monster Hunter, Resident Evil"),
        ("🕹️", "Nintendo TOKYO",                  "Única loja Nintendo do Japão — exclusivos, merchandise e edições limitadas"),
        ("⚡", "Pokemon Center Shibuya",           "Mais uma oportunidade para aquela card rara que faltou!"),
        ("🌆", "Shibuya Sky",                     "Terraço aberto no topo do Scramble Square — pôr do sol sobre Tokyo"),
        ("🚸", "Shibuya Scramble Crossing",        "O cruzamento mais famoso do mundo — foto obrigatória!"),
    ],
    "25/04 (SEX)  —  SHIBUYA • CAPCOM • NINTENDO • POKEMON • SHIBUYA SKY",
    "Capturando todos os itens raros de Tokyo!",
    "Crédito no cartão? Sim. Arrependimento? Nunca!",
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 26/04: Bate-volta Kyoto
# ══════════════════════════════════════════════════════════════════════════════
s = blank()
bg(s, RGBColor(0x1A, 0x2A, 0x1A))
for i in range(10):
    c = RGBColor(0x1A + i * 4, 0x2A + i * 3, 0x1A + i * 2)
    rect(s, 0, Inches(i * 7.5 / 10), W, Inches(7.5 / 10 + 0.05), c)

rect(s, 0, 0, W, Inches(1.45), RGBColor(0x2E, 0x7D, 0x32))
rect(s, 0, Inches(1.45), W, Inches(0.07), RGBColor(0xA5, 0xD6, 0xA7))
rect(s, 0, 0, Inches(0.4), H, RGBColor(0xA5, 0xD6, 0xA7))

circle(s, W - Inches(1.1), Inches(0.72), Inches(0.52), RGBColor(0xA5, 0xD6, 0xA7))
circle(s, W - Inches(1.1), Inches(0.72), Inches(0.36), RGBColor(0x1A, 0x2A, 0x1A))
circle(s, W - Inches(1.1), Inches(0.72), Inches(0.17), RGBColor(0xA5, 0xD6, 0xA7))

text(s, "26/04 (SAB)  —  BATE-VOLTA KYOTO", Inches(0.62), Inches(0.07),
     Inches(10.5), Inches(0.65), size=27, bold=True,
     color=RGBColor(0xA5, 0xD6, 0xA7))
text(s, "A antiga capital do Japão — templos, bambus e historia milenar",
     Inches(0.62), Inches(0.73), Inches(10.5), Inches(0.65),
     size=19, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

kyoto_items = [
    ("🚄", "Shinkansen Tokyo → Kyoto",     "~2h15 de trem bala — use o JR Pass, sai de Shinagawa ou Tokyo Station"),
    ("⛩️", "Fushimi Inari Taisha",          "Miles de portais torii vermelhos subindo a montanha — fotogênico ao maximo"),
    ("🎋", "Bosque de Bamboo (Arashiyama)", "O famoso corredor de bambus que parece outro mundo"),
    ("🏯", "Kinkaku-ji (Templo Dourado)",   "Pavilhao coberto de folhas de ouro refletido no lago"),
    ("🍵", "Chado — cerimonia do cha",      "Experience autentica de cha matcha em Gion (opcional)"),
    ("🚄", "Kyoto → Tokyo de volta",        "Retorno no fim da tarde para descansar antes da Disney!"),
]

for i, (icon, ativ, det) in enumerate(kyoto_items):
    cy = Inches(1.68 + i * 0.92)
    rect(s, Inches(0.52), cy + Inches(0.1), Inches(0.07), Inches(0.52),
         RGBColor(0x2E, 0x7D, 0x32))
    text(s, icon + "  " + ativ, Inches(0.72), cy, Inches(11.0), Inches(0.5),
         size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    text(s, det, Inches(0.72), cy + Inches(0.42), Inches(11.2), Inches(0.42),
         size=13, color=RGBColor(0xC8, 0xE6, 0xC9))

footer(s, "De trem bala ida e volta — Kyoto em um dia e voltando com a alma cheia!",
       RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xFF, 0xFF, 0xFF))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 27/04: Tokyo Disneyland
# ══════════════════════════════════════════════════════════════════════════════
make_disney_slide(
    [
        ("🎟️", "Opening → Fechamento completo",    "Aproveite cada minuto — chegue antes da abertura para pegar as filas menores"),
        ("🎢", "Space Mountain",                    "A montanha-russa no escuro — classico absoluto da Disney"),
        ("🎠", "Fantasyland",                       "O coração magico do parque: castelo, carrosséis e atrações temáticas"),
        ("🌟", "Tomorrowland & Toontown",           "Futurismo da Disney + encontro com personagens classicos"),
        ("🎭", "Shows e Parades",                   "Electrical Parade e shows ao vivo — nao perca o desfile!"),
        ("🍔", "Alimentacao no parque",             "Turkey leg, Mickey waffles e os exclusivos da Disney Japan"),
    ],
    "27/04 (DOM)  —  TOKYO DISNEYLAND",
    "Do opening ao fechamento — so Disney, nada mais!",
    "O dia mais magico da viagem — nenhum adulto envergonha de se emocionar!",
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 28/04: Tokyo DisneySea
# ══════════════════════════════════════════════════════════════════════════════
DISNEYSEA = {
    "bg":     RGBColor(0x00, 0x4E, 0x92),
    "accent": RGBColor(0x00, 0xC8, 0xFF),
    "dark":   RGBColor(0x00, 0x1E, 0x3C),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "light":  RGBColor(0xB0, 0xE8, 0xFF),
}
make_disney_slide(
    [
        ("🎟️", "Opening → Fechamento completo",  "DisneySea: o parque mais bonito do mundo (segundo os fas!)"),
        ("🌊", "Mediterranean Harbor",            "A entrada mais fotogenica de qualquer parque de diversoes"),
        ("🗺️", "Indiana Jones Adventure",         "Aventura classica em templo misterioso — fila vale a pena"),
        ("🐟", "Nemo & Friends SeaRider",         "Imersao subaquatica com projecoes — tecnologia incrivel"),
        ("🏰", "Fantasy Springs (novo!)",          "Area nova de 2024: Frozen, Rapunzel e Peter Pan — exclusivo"),
        ("🎆", "Show noturno Believe! Sea of Dreams", "Fogos + projecoes no Monte Prometheus — encerramento epico"),
    ],
    "28/04 (SEG)  —  TOKYO DISNEYSEA",
    "Do opening ao fechamento — so DisneySea, nada mais!",
    "O parque mais premiado do mundo — guarde cada segundo na memoria!",
    p=DISNEYSEA,
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 29/04: Akihabara • Rolê livre
# ══════════════════════════════════════════════════════════════════════════════
make_naruto_slide(
    [
        ("🎮", "Akihabara Electric Town",     "Paraiso otaku: games retro, figures, mangás e eletrônicos"),
        ("🃏", "Yodobashi Akiba / Mandarake", "Busca pelos colecionaveis que faltaram — ultima chance!"),
        ("🕹️", "Arcades de vários andares",   "Taiko no Tatsujin, maquinas de claw, UFO catchers e mais"),
        ("👘", "Maid Cafés (opcional)",        "Experiencia cultural única — maids cantando e servindo cafe"),
        ("🗓️", "Rolê livre em grupo",         "Qualquer canto que todo mundo queira visitar — vocês decidem!"),
        ("🛍️", "Compras de ultima hora",      "Don Quijote, combinis, souvenirs — ultima chance de encher a mala"),
    ],
    "29/04 (TER)  —  AKIHABARA & ROLEE LIVRE",
    "Missao livre — cada ninja escolhe seu caminho!",
    "Ultimo dia em Tokyo — absorva cada segundo, dattebayo!",
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 30/04: Monte Fuji • Fuji-Q • Osaka
# ══════════════════════════════════════════════════════════════════════════════
s = blank()
bg(s, FUJI["dark"])
for i in range(12):
    c = RGBColor(0x00 + i * 2, 0x2B + i * 1, 0x5C - i * 2)
    rect(s, 0, Inches(i * 7.5 / 12), W, Inches(7.5 / 12 + 0.05), c)

rect(s, 0, 0, W, Inches(1.45), FUJI["bg"])
rect(s, 0, Inches(1.45), W, Inches(0.07), FUJI["accent"])
rect(s, 0, 0, Inches(0.4), H, FUJI["accent"])

# Silhueta simplificada do Fuji (triangulo com círculos)
circle(s, W - Inches(1.3), Inches(0.72), Inches(0.58), FUJI["accent"])
circle(s, W - Inches(1.3), Inches(0.72), Inches(0.42), FUJI["dark"])
circle(s, W - Inches(1.3), Inches(0.72), Inches(0.2), FUJI["light"])

text(s, "30/04 (QUA)  —  MONTE FUJI • FUJI-Q HIGHLAND  → OSAKA",
     Inches(0.62), Inches(0.07), Inches(10.5), Inches(0.65),
     size=24, bold=True, color=FUJI["light"])
text(s, "O ponto mais alto do Japão — e depois, velocidade maxima nas montanhas-russas!",
     Inches(0.62), Inches(0.73), Inches(10.5), Inches(0.65),
     size=17, bold=True, color=FUJI["white"])

fuji_items = [
    ("🗻", "Monte Fuji — 5a Estação",       "Suba de onibus ate 2.305m — vista gelada e majestosa do vulcao"),
    ("🎢", "Fuji-Q Highland",               "Parque de diversoes com recordes mundiais: Fujiyama, Eejanaika e Do-Dodonpa"),
    ("😱", "Montanhas-russas de recorde",   "Algumas das mais radicais do mundo — estomago de aco obrigatorio"),
    ("📸", "Vista do Fuji no parque",       "Foto com o Fuji ao fundo nas atrações — impossivel de nao tirar"),
    ("🚄", "Shinkansen → Osaka",            "Monte Fuji para Osaka: use o JR Pass no Nozomi ou Hikari"),
    ("🌃", "Chegada em Osaka",              "Check-in no hotel — amanha comeca uma nova aventura na capital do street food!"),
]

for i, (icon, ativ, det) in enumerate(fuji_items):
    cy = Inches(1.68 + i * 0.92)
    rect(s, Inches(0.52), cy + Inches(0.1), Inches(0.07), Inches(0.52), FUJI["accent"])
    text(s, icon + "  " + ativ, Inches(0.72), cy, Inches(11.0), Inches(0.5),
         size=16, bold=True, color=FUJI["white"])
    text(s, det, Inches(0.72), cy + Inches(0.42), Inches(11.2), Inches(0.42),
         size=13, color=FUJI["light"])

footer(s, "Sayonara, Tokyo!  •  Proximo destino: OSAKA  🚄",
       FUJI["accent"], FUJI["dark"])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DICAS PRATICAS
# ══════════════════════════════════════════════════════════════════════════════
s = blank()
bg(s, POKEMON["dark"])
rect(s, 0, 0, W, Inches(1.3), POKEMON["bg"])
rect(s, 0, Inches(1.3), W, Inches(0.07), POKEMON["accent"])
text(s, "⚡  DICAS ESSENCIAIS — Treinador Preparado!",
     0, Inches(0.2), W, Inches(0.9), size=34, bold=True,
     color=POKEMON["accent"], align=PP_ALIGN.CENTER)

dicas = [
    ("💴", "Dinheiro",   "ATMs do 7-Eleven e Japão Post aceitam cartao estrangeiro — saque ienes"),
    ("📶", "Internet",   "Pocket Wi-Fi no aeroporto ou chip local (IIJmio) — essencial para o Maps"),
    ("🚇", "Suica Card", "Cartao recarregavel para trem, metro, onibus e combinis — vida mais facil"),
    ("🎟️", "JR Pass",   "7 dias de trem livre — vale muito para Kyoto e Osaka no final"),
    ("👟", "Calcado",    "Tenis confortavel — 15.000+ passos por dia sao garantidos"),
    ("🎒", "Bagagem",    "Servico Takuhaibin: mande as malas para o hotel por ¥1.500 — liberdade total"),
    ("📱", "Apps",       "Google Maps, Google Translate (camera), Hyperdia, Tokyo Subway Navigation"),
    ("🏥", "Saude",      "Leve seguro viagem ativo — farmácias japonesas (薬局 Yakkyoku) sao excelentes"),
]

for idx, (icon, titulo, desc) in enumerate(dicas):
    col = idx % 2
    row = idx // 2
    cx = Inches(0.4 + col * 6.5)
    cy = Inches(1.55 + row * 1.42)
    rect(s, cx, cy, Inches(0.58), Inches(0.58), POKEMON["accent"])
    text(s, icon, cx, cy, Inches(0.58), Inches(0.58),
         size=22, align=PP_ALIGN.CENTER, color=POKEMON["dark"])
    text(s, titulo, cx + Inches(0.7), cy - Inches(0.02), Inches(5.5), Inches(0.44),
         size=16, bold=True, color=POKEMON["bg"])
    text(s, desc, cx + Inches(0.7), cy + Inches(0.4), Inches(5.5), Inches(0.78),
         size=13, color=POKEMON["light"])

footer(s, "Preparado como Ash antes de entrar na Liga Pokemon!",
       POKEMON["accent"], POKEMON["dark"])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ENCERRAMENTO
# ══════════════════════════════════════════════════════════════════════════════
s = blank()
bg(s, TOKYO["dark"])
for i in range(10):
    c = RGBColor(0x0A + i * 3, 0x05 + i, 0x1A + i * 6)
    rect(s, 0, Inches(i * 0.75), W, Inches(0.8), c)

rect(s, 0, 0, Inches(0.42), H, POKEMON["accent"])
rect(s, Inches(0.42), 0, Inches(0.25), H, NARUTO["bg"])
rect(s, Inches(0.67), 0, Inches(0.25), H, JJK["accent"])
rect(s, Inches(0.92), 0, Inches(0.25), H, DISNEY["bg"])
rect(s, Inches(1.17), 0, Inches(0.25), H, FUJI["accent"])

circle(s, W - Inches(1.5), Inches(1.5), Inches(1.0), POKEMON["accent"])
circle(s, W - Inches(1.5), Inches(1.5), Inches(0.72), NARUTO["bg"])
circle(s, W - Inches(1.5), Inches(1.5), Inches(0.46), JJK["accent"])
circle(s, W - Inches(1.5), Inches(1.5), Inches(0.22), TOKYO["white"])

text(s, "SAYONARA &", Inches(1.6), Inches(1.0), Inches(9), Inches(1.0),
     size=48, bold=True, color=TOKYO["light"])
text(s, "ARIGATOU,", Inches(1.6), Inches(1.9), Inches(9), Inches(1.0),
     size=60, bold=True, color=TOKYO["accent"])
text(s, "TOKYO!  🗼", Inches(1.6), Inches(2.85), Inches(9), Inches(1.2),
     size=70, bold=True, color=TOKYO["white"])

for i, (icon, frase, cor) in enumerate([
    ("⚡", '"Gotta catch \'em all!"  — 9 dias de aventura Pokemon', POKEMON["bg"]),
    ("🍥", '"Acredite no seu jeito ninja de viajar!"  — Naruto',    NARUTO["bg"]),
    ("🩸", '"Domínio Expandido: Memorias eternas de Tokyo"  — JJK', JJK["light"]),
    ("🏰", '"The most magical trip on Earth"  — Disney + Fuji',      DISNEY["accent"]),
]):
    text(s, icon + "  " + frase, Inches(1.6), Inches(4.35 + i * 0.6),
         Inches(10.5), Inches(0.56), size=15, italic=True, color=cor)

rect(s, 0, H - Inches(0.45), W, Inches(0.45), TOKYO["accent"])
text(s, "22/04 → 30/04  •  Uma aventura que comeca com planejamento e termina com memorias para a vida  ✈️",
     0, H - Inches(0.45), W, Inches(0.45), size=13, bold=True,
     color=TOKYO["dark"], align=PP_ALIGN.CENTER)


# ── Salvar ────────────────────────────────────────────────────────────────────
output = r"c:\Users\weslley.a.oliveira\Documents\Tokyo\Roteiro_Tokyo_Anime.pptx"
prs.save(output)
print(f"Apresentacao salva: {output}")
print(f"Total de slides: {len(prs.slides)}")
